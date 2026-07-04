"""IBM-branded opportunity deck generator (python-pptx).

For each high-confidence insight, produce a short, client-ready opportunity deck
grounded in IBM brand standards (Carbon palette + IBM Plex typography):

  1. Title slide        — opportunity headline + account
  2. The 3 Why's        — Why anything / Why now / Why IBM
  3. CEM positioning    — where this sits in the Client Engagement Model
  4. MEDDPICC scorecard — qualified vs. open fields
  5. Recommended next move

Styling references IBM Carbon design tokens. If you have IBM's official
template (.potx / IBM-Presentation_Deck-Templates), you can later load it as a
base; the colour/typography choices here are aligned to that system.
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# IBM Carbon palette
IBM_BLUE = RGBColor(0x0F, 0x62, 0xFE)
IBM_DARK = RGBColor(0x16, 0x16, 0x16)
IBM_GRAY = RGBColor(0x52, 0x52, 0x52)
IBM_LIGHT = RGBColor(0xF4, 0xF4, 0xF4)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x19, 0x80, 0x38)
AMBER = RGBColor(0x8A, 0x6D, 0x00)
FONT = "IBM Plex Sans"

ACCOUNT_LABELS = {"JDI": "J.D. Irving, Limited", "IOL": "Irving Oil",
                  "GNB": "Government of New Brunswick"}

# 16:9
SW, SH = Inches(13.333), Inches(7.5)


def _slug(text: str, maxlen: int = 40) -> str:
    s = re.sub(r"[^a-zA-Z0-9]+", "-", (text or "insight")).strip("-").lower()
    return s[:maxlen] or "insight"


def _bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _textbox(slide, left, top, width, height):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    return tb, tf


def _set(run, text, size, color, bold=False, italic=False):
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic


def _accent_bar(slide, color=IBM_BLUE, top=Inches(0)):
    bar = slide.shapes.add_shape(1, Inches(0), top, Inches(0.18), SH)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()


def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


# ---------------------------------------------------------------- slides
def _title_slide(prs, ins):
    s = _blank(prs)
    _bg(s, IBM_DARK)
    _accent_bar(s, IBM_BLUE)
    _, tf = _textbox(s, Inches(0.8), Inches(2.2), Inches(11.5), Inches(3))
    p = tf.paragraphs[0]
    _set(p.add_run(), "IBM HORIZON ATLANTIC  •  OPPORTUNITY BRIEF", 13, RGBColor(0xA8, 0xA8, 0xA8))
    p2 = tf.add_paragraph()
    _set(p2.add_run(), ins.get("headline", "Opportunity Insight"), 34, WHITE, bold=True)
    p2.space_before = Pt(12)
    p3 = tf.add_paragraph()
    _set(p3.add_run(), ACCOUNT_LABELS.get(ins.get("account"), ins.get("account", "")),
         20, IBM_BLUE, bold=True)
    p3.space_before = Pt(8)
    p4 = tf.add_paragraph()
    platform = ins.get("product_platform", "")
    plat_str = f"{platform} platform   |   " if platform else ""
    _set(p4.add_run(),
         f"{plat_str}Offering: {ins.get('ibm_offering', '')}   |   "
         f"Confidence: {ins.get('confidence', 'n/a')}",
         13, RGBColor(0xC6, 0xC6, 0xC6))
    p4.space_before = Pt(16)


def _header(slide, title):
    _bg(slide, WHITE)
    _accent_bar(slide)
    _, tf = _textbox(slide, Inches(0.6), Inches(0.4), Inches(12), Inches(0.9))
    _set(tf.paragraphs[0].add_run(), title, 26, IBM_DARK, bold=True)


def _whys_slide(prs, ins):
    s = _blank(prs)
    _header(s, "The 3 Why's")
    whys = ins.get("three_whys", {})
    items = [("Why anything?", whys.get("why_anything", "")),
             ("Why now?", whys.get("why_now", "")),
             ("Why IBM?", whys.get("why_ibm", ""))]
    top = Inches(1.6)
    for label, body in items:
        card = s.shapes.add_shape(1, Inches(0.6), top, Inches(12.1), Inches(1.55))
        card.fill.solid()
        card.fill.fore_color.rgb = IBM_LIGHT
        card.line.color.rgb = IBM_BLUE
        card.line.width = Pt(0.75)
        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.15)
        _set(tf.paragraphs[0].add_run(), label, 16, IBM_BLUE, bold=True)
        p = tf.add_paragraph()
        _set(p.add_run(), body, 14, IBM_DARK)
        top += Inches(1.75)


def _cem_slide(prs, ins):
    s = _blank(prs)
    _header(s, "Client Engagement Model — positioning")
    stages = ["Prepare", "Engage", "Qualify", "Design", "Propose", "Negotiate", "Closing"]
    current = (ins.get("cem_stage", "") or "").lower()
    box_w = Inches(1.65)
    left = Inches(0.6)
    top = Inches(2.4)
    for st in stages:
        active = st.lower() in current
        box = s.shapes.add_shape(1, left, top, box_w, Inches(0.9))
        box.fill.solid()
        box.fill.fore_color.rgb = IBM_BLUE if active else IBM_LIGHT
        box.line.color.rgb = IBM_BLUE
        tf = box.text_frame
        tf.word_wrap = True
        para = tf.paragraphs[0]
        para.alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        _set(para.add_run(), st, 12, WHITE if active else IBM_GRAY, bold=active)
        left += box_w + Inches(0.08)

    _, tf = _textbox(s, Inches(0.6), Inches(3.7), Inches(12.1), Inches(2.5))
    _set(tf.paragraphs[0].add_run(), f"Stage implied: {ins.get('cem_stage', '')}",
         16, IBM_DARK, bold=True)
    p = tf.add_paragraph()
    p.space_before = Pt(10)
    _set(p.add_run(), f"Likely buyer / champion:  {ins.get('likely_buyer', 'to validate')}",
         14, IBM_DARK)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(6)
    _set(p2.add_run(), f"Trigger event:  {ins.get('trigger_type', 'n/a')}", 14, IBM_GRAY)


def _meddpicc_slide(prs, ins):
    s = _blank(prs)
    _header(s, "MEDDPICC qualification")
    medd = ins.get("meddpicc", {})
    qualified = medd.get("qualified", [])
    open_fields = medd.get("open", [])

    # Two columns
    cols = [("Qualified", qualified, GREEN), ("Still open", open_fields, AMBER)]
    left = Inches(0.6)
    for title, items, color in cols:
        _, tf = _textbox(s, left, Inches(1.6), Inches(5.9), Inches(5))
        _set(tf.paragraphs[0].add_run(), title, 18, color, bold=True)
        if items:
            for it in items:
                p = tf.add_paragraph()
                p.space_before = Pt(8)
                _set(p.add_run(), f"•  {it}", 14, IBM_DARK)
        else:
            p = tf.add_paragraph()
            p.space_before = Pt(8)
            _set(p.add_run(), "— none yet —", 13, IBM_GRAY, italic=True)
        left += Inches(6.2)


def _next_move_slide(prs, ins):
    s = _blank(prs)
    _bg(s, IBM_BLUE)
    _, tf = _textbox(s, Inches(0.9), Inches(2.3), Inches(11.5), Inches(3))
    _set(tf.paragraphs[0].add_run(), "RECOMMENDED NEXT MOVE", 14, RGBColor(0xD0, 0xE2, 0xFF))
    p = tf.add_paragraph()
    p.space_before = Pt(14)
    _set(p.add_run(), ins.get("next_move", ""), 26, WHITE, bold=True)
    p2 = tf.add_paragraph()
    p2.space_before = Pt(20)
    _set(p2.add_run(), ins.get("fact_vs_inference", ""), 12,
         RGBColor(0xD0, 0xE2, 0xFF), italic=True)
    p3 = tf.add_paragraph()
    p3.space_before = Pt(10)
    _set(p3.add_run(), f"Source: {ins.get('source_title', '')}", 11, RGBColor(0xD0, 0xE2, 0xFF))


def build_deck(ins: dict, out_dir: Path) -> Path:
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH
    _title_slide(prs, ins)
    _whys_slide(prs, ins)
    _cem_slide(prs, ins)
    _meddpicc_slide(prs, ins)
    _next_move_slide(prs, ins)

    out_dir.mkdir(parents=True, exist_ok=True)
    path = out_dir / f"opportunity_{ins.get('account', 'x')}_{_slug(ins.get('headline'))}.pptx"
    prs.save(path)
    return path


_CONF_RANK = {"low": 1, "medium": 2, "high": 3}


def build_decks(insights: list[dict], out_dir: Path,
                min_confidence: str = "medium", max_decks: int = 3) -> list[Path]:
    threshold = _CONF_RANK.get(min_confidence.lower(), 2)
    eligible = [i for i in insights
                if _CONF_RANK.get((i.get("confidence") or "").lower(), 0) >= threshold]
    paths = [build_deck(i, out_dir) for i in eligible[:max_decks]]
    return paths
